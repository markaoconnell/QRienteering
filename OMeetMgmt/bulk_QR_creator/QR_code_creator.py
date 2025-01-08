##
##  Copy & Paste Tool for images to PowerPoint(.pptx)
##
import pptx
import pptx.util
import urllib
import tempfile
import re
import sys
import getopt
from datetime import datetime

OUTPUT_TAG = "MY_TAG"
MANAGE_EVENTS = "OMeetMgmt/manage_events.php"
QR_CODE = "OMeetMgmt/create_single_qr_code.php"
REACH_CONTROL = "OMeet/reach_control.php"

TEMP_FILE_NAME = "local_temp_file.gif"

POWERPOINT_SUFFIXES = [ ".pptx", ".ppt" ]


def usage(program_name):
    print(f"Usage: {program_name}")
    print(f"{program_name}: [-h] -p site_prefix -k website access key -f output_file control_list")
    print("-h: Print this help message")
    print("-p: Prefix used to contact the site - e.g. http://www.site.com/MyOrienteeringClub/QRienteering");
    print("-k: Key used for the web site, from the administrator")
    print("-f: Output file name - e.g. QR_codes_for_event.pptx")
    print()
    print(f"Example:\n\t{program_name} -p http://www.site.com/QRienteering -k 2024_Meets -f outfile.pptx 101 102 108 109 115")


debug = True
verbose = False

site_prefix = None
site_key = None
output_file = None
output_file_suffix = ".pptx"

print(f"Starting program: {datetime.now()}");



try:
    opts, args = getopt.getopt(sys.argv[1:], "hp:f:k:")
except getopt.GetoptError:
  print("Parse error on command line.")
  usage(sys.argv[0])
  sys.exit(2)
#print "Found program arguments: ", opts
for opt, arg in opts:
  if opt == "-h":
    usage(sys.argv[0])
    sys.exit()
  elif opt == "-p":
    site_prefix = arg
  elif opt == "-f":
    output_file = arg
  elif opt == "-k":
    site_key = arg
  else:
    print (f"ERROR: Unknown option {opt}.")
    usage(sys.argv[0])
    sys.exit(1)

#if debug:
#  print("Debug is enabled.")

#if verbose:
#  print("Verbose is enabled.")


if (site_prefix == None) or (output_file == None) or (site_key == None):
    print(f"ERROR: All of -p, -k, and -f options are required and at least one is missing.")
    usage(sys.argv[0])
    sys.exit(1)

for suffix in POWERPOINT_SUFFIXES:
    if (output_file.endswith(suffix)):
        output_file_suffix = ""
        break



print(f"Parsed options, checking URL: {datetime.now()}");
try:
    my_testing_url = f"{site_prefix}/{MANAGE_EVENTS}?key={site_key}"
    with urllib.request.urlopen(my_testing_url) as response:
       testing_response = response.read()
except urllib.error.URLError:
    print(f"Error attempting to query URL - is it valid?  Please confirm in a browser.")
    print(f"URL: {my_testing_url}")
    sys.exit(1)


decoded_response = testing_response.decode("utf-8")
#print(f"Got {decoded_response} in repsonse to query of {my_testing_url}")
found_new_event_create = re.search(r"####,XLATED_KEY,", decoded_response)
if found_new_event_create == None:
    print(f"Invalid response received querying URL - is it valid?  Please confirm in a browser.")
    print(f"URL: {my_testing_url}")
    print(f"Response: {decoded_response}")
    sys.exit(1)

   
print(f"Validated URL, creating presentation: {datetime.now()}");

# new
prs = pptx.Presentation()
# open
# prs_exists = pptx.Presentation("some_presentation.pptx")

# default slide width
#prs.slide_width = 9144000
# slide height @ 4:3
#prs.slide_height = 6858000
# slide height @ 16:9
#prs.slide_height = 5143500

prs.slide_height = pptx.util.Inches(9)
prs.slide_width = pptx.util.Inches(12)

#prs.slide_height = pptx.util.Inches(16)
#prs.slide_width = pptx.util.Inches(9)

# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
# blank slide
#slide = prs.slides.add_slide(prs.slide_layouts[6])

# set title
title = slide.shapes.title
title.text = f"QR codes for site: {site_prefix}"

pic_left  = int(prs.slide_width * 0.15)
pic_top   = int(prs.slide_height * 0.1)
pic_width = int(prs.slide_width * 0.7)




#print(f"Got {qrcode_image_bytes} in repsonse to query of {my_qrcode_url}")

for control_id in args:
    print (f"Working on {control_id}: {datetime.now()}\n")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

#    tb = slide.shapes.add_textbox(-int(prs.slide_height * 0.4), int(prs.slide_height * 0.4), prs.slide_height, int(prs.slide_width * 0.03))
    tb = slide.shapes.add_textbox(-pptx.util.Inches(4.25), pptx.util.Inches(4.25), prs.slide_height, pptx.util.Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = pptx.enum.text.PP_ALIGN.CENTER
    tb.rotation = -90
    p.text = f"{control_id}"
    p.font.size = pptx.util.Pt(138)
    p.font.underline = True

    #pic_height = int(pic_width * img.shape[0] / img.shape[1])
    #pic   = slide.shapes.add_picture(g, pic_left, pic_top)

    my_qrcode_url = f"{site_prefix}/{QR_CODE}?qr_code={site_prefix}/{REACH_CONTROL}?control={control_id}"
    with urllib.request.urlopen(my_qrcode_url) as response:
       qrcode_image_bytes = response.read()

    # it seems to come in with a leading newline, strip it off
    qrcode_image_bytes = qrcode_image_bytes.lstrip()

    with open(TEMP_FILE_NAME, "wb") as qrcode_file:
        qrcode_file.write(qrcode_image_bytes)

    pic_left = pptx.util.Inches(2.5)
    pic_top = pptx.util.Inches(1.25)
    pic_height = pptx.util.Inches(6.5)
    pic_width = pic_height

    pic   = slide.shapes.add_picture(TEMP_FILE_NAME, pic_left, pic_top, pic_width, pic_height)

    #pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)


print (f"Saving presentation: {datetime.now()}\n")
prs.save(f"{output_file}{output_file_suffix}")
print (f"All done: {datetime.now()}\n")
